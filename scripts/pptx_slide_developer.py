from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]  # Blank layout

# === Slide 1: Title Slide with image on left, text on right ===
slide1 = prs.slides.add_slide(blank_slide_layout)

# Add image on the left
img_path1 = 'title_image.jpg'
slide1.shapes.add_picture(img_path1, Inches(0.5), Inches(1), width=Inches(4.5))

# Add text box on the right
left = Inches(5.2)
top = Inches(1)
width = Inches(4)
height = Inches(3)
textbox1 = slide1.shapes.add_textbox(left, top, width, height)
text_frame1 = textbox1.text_frame
text_frame1.text = "Welcome to the Presentation"

p = text_frame1.add_paragraph()
p.text = "Generated using Python"
p.font.size = Pt(20)
p.font.bold = True

# === Slide 2: Plot Slide with text on left, image on right, and hyperlink ===
slide2 = prs.slides.add_slide(blank_slide_layout)

# Text box on the left
textbox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1), Inches(4), Inches(3))
text_frame2 = textbox2.text_frame
text_frame2.text = "Here's a static preview of the plot."

p2 = text_frame2.add_paragraph()
p2.text = "Click below to view the interactive version."
p2.space_before = Pt(14)

# Add hyperlink shape (simple rectangle)
left = Inches(0.5)
top = Inches(3.5)
width = Inches(2.5)
height = Inches(0.6)
link_box = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
link_box.text = "View Interactive Plot"

# Format and add hyperlink
fill = link_box.fill
fill.solid()
fill.fore_color.rgb = RGBColor(91, 155, 213)
link_box.text_frame.paragraphs[0].font.size = Pt(14)
link_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
link_box.click_action.hyperlink.address = "https://example.com/interactive_plot.html"

# Add plot image on the right
img_path2 = 'plot.jpg'
slide2.shapes.add_picture(img_path2, Inches(5.2), Inches(1), width=Inches(4))

# === Slide 3: Table Slide ===
slide3 = prs.slides.add_slide(blank_slide_layout)

# Table parameters
rows, cols = 5, 3
table_data = [
    ["Metric", "Value", "Comment"],
    ["Sales", "1000", "Q1 Sales"],
    ["Revenue", "2500", "Q1 Revenue"],
    ["Profit", "400", "Q1 Profit"],
    ["Margin", "16%", "Q1 Margin"]
]

left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(3)
table = slide3.shapes.add_table(rows, cols, left, top, width, height).table

# Fill table with data
for r in range(rows):
    for c in range(cols):
        cell = table.cell(r, c)
        cell.text = table_data[r][c]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)

# Save the presentation
prs.save("\visualizations\custom_presentation.pptx")